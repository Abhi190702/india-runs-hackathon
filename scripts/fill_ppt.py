# -*- coding: utf-8 -*-
"""Fill the India.Runs Idea Submission template with RedrobRank V2 content."""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR

SRC = r"data/Copy of Idea Submission Template _ Redrob.pptx"
OUT = r"data/RedrobRank_V2_Submission_FILLED.pptx"

prs = Presentation(SRC)


def shape_by_id(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    return None


def base_color(shape):
    try:
        c = shape.text_frame.paragraphs[0].runs[0].font.color
        if c and c.type is not None:
            return c.rgb
    except Exception:
        pass
    return RGBColor(0x20, 0x20, 0x20)


def fill(shape, items, size=14):
    """items: list of (text, level, bold)."""
    color = base_color(shape)
    tf = shape.text_frame
    tf.word_wrap = True
    try:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except Exception:
        pass
    tf.clear()
    for i, (text, level, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        try:
            r.font.color.rgb = color
        except Exception:
            pass


# ---- Slide 0: Problem Statement (append value to the existing label) --------
s0 = prs.slides[0]
ps = shape_by_id(s0, 56)
stmt = ("Intelligent Candidate Discovery & Ranking - rank the top 100 of "
        "100,000 candidate profiles for a Senior AI Engineer (Founding Team) "
        "role at Redrob AI by understanding who truly fits, not by matching keywords.")
p = ps.text_frame.paragraphs[0]
src_font = p.runs[0].font if p.runs else None
run = p.add_run()
run.text = stmt
run.font.size = Pt(15)
if src_font is not None:
    if src_font.name:
        run.font.name = src_font.name
    try:
        if src_font.color and src_font.color.type is not None:
            run.font.color.rgb = src_font.color.rgb
    except Exception:
        pass

# ---- Slide 1: Solution Overview --------------------------------------------
fill(shape_by_id(prs.slides[1], 64), [
    ("Proposed solution", 0, True),
    ("RedrobRank V2 - a deterministic, CPU-only, fully explainable ranking engine.", 1, False),
    ("Reads the JD, scores every candidate 0-1, returns the top 100 with a written reason for each.", 1, False),
    ('Core philosophy: "Career proof beats keyword claims."', 1, False),
    ("What makes it different", 0, True),
    ("Source-aware evidence: a skill in job history counts far more (1.00) than in a skills list (0.25) - defeats keyword-stuffers.", 1, False),
    ("Honeypot / anomaly auditor: pure logic catches impossible profiles that embeddings cannot.", 1, False),
    ("Hireability signals: down-weights unreachable candidates (inactive, unresponsive, long notice).", 1, False),
    ("No black box: every score traces to inspectable evidence - built to be defended.", 1, False),
], size=13)

# ---- Slide 2: JD Understanding ---------------------------------------------
fill(shape_by_id(prs.slides[2], 71), [
    ("Key requirements extracted from the JD", 0, True),
    ("Production retrieval / ranking / search / recommendation / embeddings experience.", 1, False),
    ("Vector & hybrid search (FAISS, BM25, Pinecone); evaluation (NDCG / MRR / MAP).", 1, False),
    ("5-9 yrs, product-company background; NOT keyword-stuffers, pure-research, consulting-only, or title-chasers.", 1, False),
    ("Evaluating fit beyond keyword matching", 0, True),
    ('We don\'t ask "does the word appear?" - we ask "where did the evidence come from, and does the work history prove it?"', 1, False),
    ("Top signals: career-backed retrieval/ranking evidence, production-systems evidence, experience band, hireability.", 1, False),
    ("Plain-language gems (e.g. a Data Scientist who built a recommender) are surfaced even without buzzwords.", 1, False),
], size=13)

# ---- Slide 3: Ranking Methodology ------------------------------------------
fill(shape_by_id(prs.slides[3], 78), [
    ("Retrieve -> score -> rank", 0, True),
    ("Parse 100k profiles -> source-aware evidence -> normalized features -> local semantic signal -> anomaly penalties & caps -> weighted fusion -> deterministic sort -> top 100.", 1, False),
    ("Models / algorithms / heuristics", 0, True),
    ("Rule-driven, inspectable scoring (explainable by design).", 1, False),
    ("Local TF-IDF semantic similarity (scikit-learn) - no network, no hosted LLM.", 1, False),
    ("Logic-based honeypot / consistency auditor.", 1, False),
    ("Combining signals into the final score", 0, True),
    ("final = clamp( min( 0.72*technical + 0.18*hireability + 0.10*trust - anomaly_penalty , score_cap ), 0, 1 )", 1, False),
    ("Deterministic tie-breaks (down to candidate_id) for reproducibility.", 1, False),
], size=13)

# ---- Slide 4: Explainability & Data Validation -----------------------------
fill(shape_by_id(prs.slides[4], 85), [
    ("How ranking decisions are explained", 0, True),
    ("Each top-100 candidate gets a 1-2 sentence reason citing specific profile facts and specific JD requirements (and honest concerns at lower ranks).", 1, False),
    ("Preventing hallucinations / unsupported claims", 0, True),
    ("Reasoning is assembled deterministically from extracted facts and quotes the candidate's own career text - it cannot invent skills or employers. 100/100 reasonings unique.", 1, False),
    ("Handling inconsistent / low-quality / suspicious profiles", 0, True),
    ("Hard honeypots (impossible profiles, e.g. 8 yrs at a 3-yr-old company) -> detected by logic and floored.", 1, False),
    ("Soft anomalies (keyword-stuffing, shallow-LLM-only, consulting-only) -> penalties.", 1, False),
    ("Score caps -> suspicious profiles can't reach the top. Result: 0 honeypots in the top 100.", 1, False),
], size=12)

# ---- Slide 5: End-to-End Workflow ------------------------------------------
fill(shape_by_id(prs.slides[5], 92), [
    ("Job Description", 0, False),
    ("-> JD rubric (must-haves, disqualifiers, weights)", 1, False),
    ("-> Parse 100k candidate profiles", 1, False),
    ("-> Source-aware evidence extraction", 1, False),
    ("-> Features: technical / hireability / trust (each 0-1)", 1, False),
    ("-> Local TF-IDF semantic boost", 1, False),
    ("-> Anomaly + honeypot detection -> score caps", 1, False),
    ("-> Weighted scoring + deterministic sort", 1, False),
    ("-> Top 100 + grounded reasoning -> submission.csv", 1, False),
    ("One command: python src/rank.py --candidates ./data/candidates.jsonl --out ./submission.csv", 0, True),
], size=14)

# ---- Slide 6: System Architecture (no content box -> add one) --------------
s6 = prs.slides[6]
box = s6.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.4))
arch_lines = [
    "data/candidates.jsonl",
    "    parser.py      ->  clean + flatten 100k profiles",
    "    evidence.py    ->  source-aware concept matching",
    "    features.py    ->  technical / hireability / trust features",
    "    semantic.py    ->  local TF-IDF meaning match (no network)",
    "    anomalies.py   ->  hard honeypots + soft anomalies",
    "    caps.py        ->  ceilings for suspicious profiles",
    "    scoring.py     ->  weighted fusion + deterministic sort",
    "    reasoning.py   ->  grounded, non-hallucinated reasons",
    "    rank.py        ->  CLI entry point",
    "          |",
    "          v",
    "    submission.csv (top 100)   +   Streamlit demo sandbox",
    "",
    "Constraints honored: CPU-only . no network . no hosted-LLM at ranking . deterministic.",
]
tf = box.text_frame
tf.word_wrap = True
for i, line in enumerate(arch_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r = p.add_run()
    r.text = line
    r.font.name = "Consolas"
    r.font.size = Pt(13 if i == len(arch_lines) - 1 else 12)
    r.font.bold = (i == len(arch_lines) - 1)
    r.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

# ---- Slide 7: Results & Performance ----------------------------------------
fill(shape_by_id(prs.slides[7], 105), [
    ("Ranking quality (full 100,000-candidate run)", 0, True),
    ("Top-100 is 100% genuine AI / ML / Search / NLP / Data titles - 0 wrong-domain, 0 honeypots.", 1, False),
    ("69 honeypots detected pool-wide (about the ~80 planted).", 1, False),
    ("12 plain-language gems (non-AI-title builders) correctly surface into the top 100.", 1, False),
    ("NDCG@10 = 1.00, NDCG@50 = 0.91, P@10 = 1.00 vs our offline recruiter-aligned gold set.", 1, False),
    ("100/100 reasonings unique (no templating, no hallucination).", 1, False),
    ("Runtime & compute (budget: 5 min, 16 GB, CPU, no network)", 0, True),
    ("~55 sec for 100k (parallel); ~3.4 min fully serial - both well under budget.", 1, False),
    ("Deterministic, with a serial fallback if a sandbox blocks parallelism. 20/20 tests pass.", 1, False),
], size=13)

# ---- Slide 8: Technologies Used --------------------------------------------
fill(shape_by_id(prs.slides[8], 112), [
    ("Python 3.11 - core language; clean and reproducible.", 0, False),
    ("scikit-learn (TF-IDF) - local semantic similarity, no network / API.", 0, False),
    ("NumPy - fast vector math within the CPU budget.", 0, False),
    ("Streamlit - live demo sandbox (required deliverable).", 0, False),
    ("pytest - 20 automated tests for correctness & reproducibility.", 0, False),
    ("Git / GitHub - versioned, staged history (real iteration, not a single dump).", 0, False),
    ("Deliberately no GPU and no hosted LLM at ranking time - the system must scale to 200k+ profiles in production within tight latency, so we chose lightweight, local, explainable components.", 0, True),
], size=14)

# ---- Slide 9: Submission Assets --------------------------------------------
fill(shape_by_id(prs.slides[9], 119), [
    ("GitHub repo: github.com/hazelmayank/india-runs-hackathon", 0, False),
    ("Live demo (Streamlit): india-runs-hackathon-d4kd2fvfgaxeqbzen3hhyx.streamlit.app", 0, False),
    ("Ranked output: submission.csv (top 100, official format, validator-passed).", 0, False),
    ("Docs: README, architecture, methodology, honeypot, reproducibility notes.", 0, False),
    ("Reproduce: python src/rank.py --candidates ./data/candidates.jsonl --out ./submission.csv", 0, False),
], size=15)

prs.save(OUT)
print("Saved:", OUT)
